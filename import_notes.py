#!/usr/bin/env python3
"""Import notes from a local directory into the Obsidian vault.

Supports three directory layouts (auto-detected):

1. **Customer-first**: Top-level dirs are customer names with mixed content
   beneath.  The importer infers customer from the directory name and looks at
   subdirectory names (Meetings/, QBR/, Opportunities/, etc.) for type hints.

       Customers/
       ├── Nomura/
       │   ├── Nomura Notes.docx         <- customer=Nomura, type from AI
       │   ├── Meetings/                 <- type=meeting
       │   │   └── weekly-sync.docx
       │   └── Archived/                 <- imported, flagged status=archived
       │       └── old-doc.docx
       └── M&G/
           └── Discovery Capture Sheet.docx

2. **Type-first**: Top-level dirs match note types.  Type is set from the
   directory; optional second level is the customer.

       meeting/acme/sync.docx  ->  type=meeting, customer=Acme

3. **Unsorted / flat**: Loose files in the root -- classified by filename patterns.

Duplicate files with ``(1)`` suffixes and meta files (``_notebook-instructions.txt``)
are automatically skipped.

After import, use Cursor to review and reclassify notes for best results
(no separate API key needed).

Usage:
    python import_notes.py ~/Downloads/Customers --dry-run
    python import_notes.py ~/Downloads/Customers
"""
from __future__ import annotations

import argparse
import json
import logging
import re
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import yaml

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger(__name__)

PROJECT_ROOT = Path(__file__).resolve().parent

SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".md", ".txt", ".doc"}

TYPE_TO_FOLDER = {
    "customer": "01-Customers",
    "person": "01-Customers",
    "project": "02-Projects",
    "meeting": "03-Meetings",
    "knowledge": "04-Knowledge",
    "industry": "05-Industry",
    "competition": "06-Competition",
    "demo": "07-Demos",
    "esat": "08-ESATs",
    "partner": "11-Partners",
    "workshop": "09-Workshops",
    "certification": "10-Certifications",
    "internal": "12-Internal",
}

# Maps common subdirectory names (in a customer-first layout) to note types
SUBDIR_TYPE_HINTS: dict[str, str | None] = {
    "meetings": "meeting",
    "meeting": "meeting",
    "qbr": "meeting",
    "opportunities": "project",
    "opportunity": "project",
    "rfp": "project",
    "esat": "esat",
    "esats": "esat",
    "demos": "demo",
    "demo": "demo",
    "archived": None,
    "notebooklm": "customer",
    "open": None,
}

SKIP_FILENAMES = {"_notebook-instructions.txt", ".ds_store"}

DUPLICATE_SUFFIX = re.compile(r"\(\d+\)$")


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

@dataclass
class ImportConfig:
    vault_path: Path
    folders: list[str]
    note_types: list[str]
    tracker_path: Path


def load_config(config_path: Path) -> ImportConfig:
    """Load vault config for import."""
    with config_path.open(encoding="utf-8") as f:
        raw = yaml.safe_load(f)

    vault_path = Path(raw["vault"]["path"]).expanduser().resolve()

    return ImportConfig(
        vault_path=vault_path,
        folders=raw.get("folders", []),
        note_types=raw.get("note_types", []),
        tracker_path=PROJECT_ROOT / ".import_history.json",
    )


# ---------------------------------------------------------------------------
# Import history
# ---------------------------------------------------------------------------

def load_history(path: Path) -> dict[str, Any]:
    """Load previously imported files to avoid duplicates."""
    if path.exists():
        with path.open(encoding="utf-8") as f:
            return json.load(f)
    return {}


def save_history(path: Path, history: dict[str, Any]) -> None:
    """Persist import history."""
    with path.open("w", encoding="utf-8") as f:
        json.dump(history, f, indent=2)


# ---------------------------------------------------------------------------
# File discovery and path metadata extraction
# ---------------------------------------------------------------------------

DATE_PATTERN = re.compile(r"(\d{4}[-_]\d{2}[-_]\d{2})")


@dataclass
class SourceFile:
    path: Path
    hint_type: str | None
    hint_customer: str | None
    hint_date: str | None
    hint_archived: bool
    relative_parts: list[str]
    extra_context: list[str] = field(default_factory=list)


def _normalise_for_match(name: str) -> str:
    """Strip hyphens, underscores, spaces for fuzzy matching."""
    return name.lower().replace("-", "").replace("_", "").replace(" ", "")


def _extract_date_from_name(name: str) -> str | None:
    """Extract YYYY-MM-DD from a filename or folder name."""
    match = DATE_PATTERN.search(name)
    if match:
        return match.group(1).replace("_", "-")
    return None


def _is_duplicate(file_path: Path, all_stems: set[str]) -> bool:
    """Check if a file looks like a Google Drive duplicate, e.g. 'Report(1).docx'."""
    stem = file_path.stem
    match = DUPLICATE_SUFFIX.search(stem)
    if not match:
        return False
    base_stem = stem[:match.start()].rstrip()
    return base_stem in all_stems


def discover_files(source_dir: Path, note_types: list[str]) -> list[SourceFile]:
    """Find importable files, auto-detecting customer-first or type-first layout.

    Heuristic: if the first-level directory names mostly DON'T match known note
    types, treat the layout as customer-first. Otherwise type-first.
    """
    type_lookup = {_normalise_for_match(nt): nt for nt in note_types}

    # Collect all candidate files first (for dedup scanning)
    candidates: list[Path] = []
    for file_path in sorted(source_dir.rglob("*")):
        if not file_path.is_file():
            continue
        if file_path.name.lower() in SKIP_FILENAMES:
            continue
        if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        candidates.append(file_path)

    if not candidates:
        return []

    # Build set of base stems per directory for dedup
    dir_stems: dict[Path, set[str]] = {}
    for fp in candidates:
        dir_stems.setdefault(fp.parent, set()).add(fp.stem)

    # Auto-detect layout: check top-level directory names
    top_dirs = {
        fp.relative_to(source_dir).parts[0]
        for fp in candidates
        if len(fp.relative_to(source_dir).parts) > 1
    }
    type_matches = sum(1 for d in top_dirs if _normalise_for_match(d) in type_lookup)
    customer_first = len(top_dirs) > 0 and type_matches < len(top_dirs) / 2

    if customer_first:
        logger.info("Detected customer-first directory layout")
    else:
        logger.info("Detected type-first directory layout")

    results: list[SourceFile] = []
    dup_count = 0

    for file_path in candidates:
        # Skip Google-Drive-style duplicates: "file(1).docx"
        if _is_duplicate(file_path, dir_stems.get(file_path.parent, set())):
            dup_count += 1
            continue

        rel = file_path.relative_to(source_dir)
        parts = list(rel.parts[:-1])

        hint_type: str | None = None
        hint_customer: str | None = None
        hint_date: str | None = None
        hint_archived = False
        extra_context: list[str] = []

        if customer_first and parts:
            # First level = customer name
            hint_customer = parts[0]

            # Scan remaining directory levels for type hints and context
            for part in parts[1:]:
                norm = _normalise_for_match(part)

                if norm == "archived":
                    hint_archived = True

                if norm in SUBDIR_TYPE_HINTS and hint_type is None:
                    hint_type = SUBDIR_TYPE_HINTS[norm]

                # Deeper levels (e.g. opportunity names) add context
                if norm not in SUBDIR_TYPE_HINTS:
                    extra_context.append(part)

        elif parts:
            # Type-first layout (original behaviour)
            normalised = _normalise_for_match(parts[0])
            if normalised in type_lookup:
                hint_type = type_lookup[normalised]

                if len(parts) >= 2 and not DATE_PATTERN.match(parts[1]):
                    hint_customer = parts[1]

        # Extract date from any path component or the filename
        for part in parts:
            found = _extract_date_from_name(part)
            if found:
                hint_date = found
                break
        if not hint_date:
            hint_date = _extract_date_from_name(file_path.stem)

        results.append(SourceFile(
            path=file_path,
            hint_type=hint_type,
            hint_customer=hint_customer,
            hint_date=hint_date,
            hint_archived=hint_archived,
            relative_parts=parts,
            extra_context=extra_context,
        ))

    if dup_count:
        logger.info("Skipped %d duplicate files (Google Drive copies)", dup_count)

    return results


# ---------------------------------------------------------------------------
# Conversion
# ---------------------------------------------------------------------------

def _get_pandoc_path() -> str:
    """Locate the pandoc binary (prefers the one bundled with pypandoc_binary)."""
    try:
        import pypandoc
        return pypandoc.get_pandoc_path()
    except (ImportError, OSError):
        pass
    import shutil
    path = shutil.which("pandoc")
    if path:
        return path
    raise FileNotFoundError(
        "Pandoc not found. Run: pip install pypandoc_binary"
    )


def convert_docx_to_markdown(docx_path: Path) -> str:
    """Convert a .docx or .doc file to Markdown.

    Calls pandoc directly via subprocess to avoid pypandoc's glob expansion
    which breaks on filenames containing brackets.
    """
    pandoc = _get_pandoc_path()
    result = subprocess.run(
        [pandoc, str(docx_path), "-f", "docx", "-t", "markdown", "--wrap=none"],
        capture_output=True,
        text=True,
        check=True,
    )
    return result.stdout


def extract_pptx_text(pptx_path: Path) -> str:
    """Extract all text content and speaker notes from a .pptx file.

    Reads text from shapes (titles, bullet points, text boxes) and speaker
    notes on each slide. Cannot read text embedded in images or charts.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx not installed. Run: pip install python-pptx"
        )

    prs = Presentation(str(pptx_path))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_cells:
                        slide_texts.append(" | ".join(row_cells))

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        parts.append(f"## Slide {i}")
        if slide_texts:
            parts.append("\n".join(slide_texts))
        if notes_text:
            parts.append(f"\n**Speaker Notes:**\n{notes_text}")
        parts.append("")

    return "\n\n".join(parts)


def read_text_file(file_path: Path) -> str:
    """Read a .md or .txt file."""
    return file_path.read_text(encoding="utf-8")


def convert_file(file_path: Path) -> str:
    """Convert any supported file to Markdown content."""
    ext = file_path.suffix.lower()

    if ext in (".docx", ".doc"):
        return convert_docx_to_markdown(file_path)
    elif ext == ".pptx":
        return extract_pptx_text(file_path)
    elif ext in (".md", ".txt"):
        return read_text_file(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ---------------------------------------------------------------------------
# Local classification (no API required)
# ---------------------------------------------------------------------------

# Filename patterns that hint at a note type regardless of directory
FILENAME_TYPE_PATTERNS: list[tuple[re.Pattern[str], str]] = [
    (re.compile(r"(?i)\bqbr\b"), "meeting"),
    (re.compile(r"(?i)\baccount\s*plan\b"), "customer"),
    (re.compile(r"(?i)\bdiscovery\s*capture\b"), "customer"),
    (re.compile(r"(?i)\brenewal\b"), "customer"),
    (re.compile(r"(?i)\bweekly\s*(sync|update|meeting)"), "meeting"),
    (re.compile(r"(?i)\bsync\b.*notes"), "meeting"),
    (re.compile(r"(?i)\bnotes\s*by\s*gemini\b"), "meeting"),
    (re.compile(r"(?i)\bproposal\b"), "project"),
    (re.compile(r"(?i)\bdemo\s*brief\b"), "demo"),
    (re.compile(r"(?i)\besat\b"), "esat"),
    (re.compile(r"(?i)\bplayback\b"), "demo"),
    (re.compile(r"(?i)\bhandout\b"), "knowledge"),
    (re.compile(r"(?i)\btraining\b"), "certification"),
    (re.compile(r"(?i)\brfi\b|\brfp\b"), "project"),
    (re.compile(r"(?i)\barchitecture\b"), "knowledge"),
    (re.compile(r"(?i)\baccount\s*notes\b"), "customer"),
    (re.compile(r"(?i)\bactivity\s*log\b"), "customer"),
    (re.compile(r"(?i)\bdeal\s*history\b"), "customer"),
    (re.compile(r"(?i)\bsalesforce\b"), "customer"),
    (re.compile(r"(?i)\bsupport\s*cases\b"), "customer"),
    (re.compile(r"(?i)\bcloud\s*telemetry\b"), "customer"),
    (re.compile(r"(?i)\bweekly\s*updates?\b"), "meeting"),
]


ACRONYMS = {
    "esat", "qbr", "rfi", "rfp", "elk", "apm", "aro", "ctap", "plc",
    "sdk", "gmt", "bst", "uk", "hd", "elmo", "inc", "ko",
}

GEMINI_TIMESTAMP = re.compile(
    r"\s*[-_]?\s*\d{4}[-_ ]\d{2}[-_ ]\d{2}\s+\d{2}[-_ ]\d{2}\s*(GMT|BST|UTC)\s*",
    re.IGNORECASE,
)
GEMINI_SUFFIX = re.compile(r"\s*[-_]?\s*Notes\s+by\s+Gemini\s*$", re.IGNORECASE)
COPY_PREFIX = re.compile(r"^Copy\s+of\s+", re.IGNORECASE)


def _smart_title_case(text: str) -> str:
    """Title-case while preserving known acronyms."""
    words = text.split()
    result = []
    for word in words:
        if word.lower() in ACRONYMS:
            result.append(word.upper())
        elif word.isupper() and len(word) > 1:
            result.append(word)
        else:
            result.append(word.capitalize())
    return " ".join(result)


def _clean_title(filename_stem: str, customer: str | None) -> str:
    """Turn a raw filename stem into a clean note title."""
    title = filename_stem

    # Strip leading dates like "2019-07 " or "2021-05 "
    title = re.sub(r"^\d{4}[-_]\d{2}[-_]?\d{0,2}\s*[-_]?\s*", "", title)

    # Strip trailing version markers like "_v1.3", "-v2", " v4.1"
    title = re.sub(r"[\s_-]*v\d+(\.\d+)*$", "", title, flags=re.IGNORECASE)

    # Strip "Copy of " prefix
    title = COPY_PREFIX.sub("", title)

    # Strip Gemini timestamp junk: "2025_11_14 14_58 GMT"
    title = GEMINI_TIMESTAMP.sub("", title)

    # Strip "Notes by Gemini" suffix
    title = GEMINI_SUFFIX.sub("", title)

    # Clean up separators
    title = title.replace("_", " ")
    title = re.sub(r"\s+", " ", title).strip(" -_")

    # Always strip customer name from the title
    if customer:
        escaped = re.escape(customer)
        title = re.sub(rf"(?i)\s*[-_]?\s*{escaped}\s*[-_]?\s*", " ", title)
        title = re.sub(rf"(?i)^{escaped}\s*[-_&]\s*", "", title)
        # Clean up leading connectors left after stripping ("& Elastic" → "Elastic")
        title = re.sub(r"^[&,\s_-]+", "", title)
        title = title.strip(" -_")

    # Smart title case (preserves acronyms)
    if title:
        title = _smart_title_case(title)

    return title or _smart_title_case(filename_stem)


def _infer_type_from_filename(stem: str) -> str | None:
    """Check filename against known patterns to guess a note type."""
    for pattern, note_type in FILENAME_TYPE_PATTERNS:
        if pattern.search(stem):
            return note_type
    return None


def classify_locally(sf: SourceFile) -> dict[str, Any]:
    """Classify a file using only path and filename signals (no API calls)."""
    note_type = sf.hint_type

    # If directory didn't give us a type, try filename patterns
    if not note_type:
        note_type = _infer_type_from_filename(sf.path.stem)

    # Default: if under a customer directory, call it a "customer" note
    if not note_type and sf.hint_customer:
        note_type = "customer"

    # Last resort
    if not note_type:
        note_type = "knowledge"

    target_folder = TYPE_TO_FOLDER.get(note_type, "00-Inbox")
    title = _clean_title(sf.path.stem, sf.hint_customer)

    # Prefix customer name to disambiguate notes that would collide across customers
    if sf.hint_customer and len(title.split()) <= 3:
        title = f"{sf.hint_customer} - {title}"

    # For meetings, prefix with date
    if sf.hint_date and note_type in ("meeting",):
        title = f"{sf.hint_date} - {title}"

    frontmatter: dict[str, Any] = {"type": note_type}

    if sf.hint_customer:
        frontmatter["customer"] = f"[[{sf.hint_customer}]]"
    if sf.hint_date:
        frontmatter["date"] = sf.hint_date
    if sf.hint_archived:
        frontmatter["status"] = "archived"
    if sf.extra_context:
        frontmatter["context"] = " / ".join(sf.extra_context)

    return {
        "title": title,
        "frontmatter": frontmatter,
        "target_folder": target_folder,
        "type": note_type,
        "confidence": 0.7 if sf.hint_type else 0.4,
        "reasoning": f"Local: type={'dir' if sf.hint_type else 'filename/default'}, customer={'dir' if sf.hint_customer else 'none'}",
    }


# ---------------------------------------------------------------------------
# AI: metadata extraction (type already known)
# ---------------------------------------------------------------------------

def _build_path_context(sf: SourceFile) -> str:
    """Build a structured context block from path-derived metadata."""
    lines = [f"Original filename: {sf.path.name}"]
    if sf.relative_parts:
        lines.append(f"Directory path: {'/'.join(sf.relative_parts)}/")
    if sf.hint_customer:
        lines.append(f"Customer hint (from folder name): {sf.hint_customer}")
    if sf.hint_date:
        lines.append(f"Date hint (from path/filename): {sf.hint_date}")
    if sf.hint_archived:
        lines.append("Note: this file was in an 'Archived' directory")
    if sf.extra_context:
        lines.append(f"Additional path context: {' / '.join(sf.extra_context)}")
    return "\n".join(lines)


def extract_metadata(
    content: str,
    sf: SourceFile,
    known_type: str,
) -> dict[str, Any]:
    """Use Claude to extract metadata when the note type is already known.

    Path-derived hints (customer, date) are passed as structured context so
    Claude can incorporate them alongside content analysis.
    """
    try:
        import anthropic
    except ImportError:
        raise ImportError(
            "Anthropic SDK not installed. Run: pip install anthropic"
        )

    client = anthropic.Anthropic()
    path_context = _build_path_context(sf)

    prompt = f"""Extract metadata from this document for an Obsidian vault note.

The note type is already known: **{known_type}**

The vault belongs to a Solution Architect at Elastic.

File context (use these hints to fill metadata where the document content is ambiguous):
{path_context}

Document content (first 4000 chars):
{content[:4000]}

Respond with ONLY valid JSON (no markdown fences):
{{
  "title": "<Title Case note title, no file extension>",
  "frontmatter": {{
    "type": "{known_type}",
    ...additional relevant properties...
  }}
}}

Rules for frontmatter properties by type:
- meeting: include date (YYYY-MM-DD), customer ("[[Name]]"), attendees if identifiable
- customer: include industry, status ("active"), tier if identifiable
- person: include customer ("[[Name]]"), role
- project: include customer ("[[Name]]"), status, start-date
- knowledge: include product-area (Search, Observability, Security, Platform)
- competition: include competitor, product-area
- certification: include provider, status, target-date if identifiable
- industry: include vertical
- demo: include product, customer ("[[Name]]") if applicable
- workshop: include customer ("[[Name]]"), duration, date

Rules:
- Use wikilinks for customer/project/people references: "[[Name]]"
- Title should be Title Case with spaces
- For dated content, prefix title with YYYY-MM-DD - Title
- Only include properties that have actual values extracted from the content
- If a customer or date hint is provided above, use it unless the document clearly contradicts it"""

    message = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=1024,
        messages=[{"role": "user", "content": prompt}],
    )

    return json.loads(message.content[0].text)


# ---------------------------------------------------------------------------
# AI: full classification (type unknown)
# ---------------------------------------------------------------------------

def classify_document(
    content: str,
    sf: SourceFile,
    folders: list[str],
    note_types: list[str],
) -> dict[str, Any]:
    """Use Claude to classify type AND extract metadata when type is unknown.

    Any path-derived hints (customer, date) are still passed as context even
    when the note type couldn't be determined from the directory.
    """
    try:
        import anthropic
    except ImportError:
        raise ImportError(
            "Anthropic SDK not installed. Run: pip install anthropic"
        )

    client = anthropic.Anthropic()

    folder_list = "\n".join(f"  - {f}" for f in folders)
    type_list = ", ".join(note_types)
    path_context = _build_path_context(sf)

    prompt = f"""Analyse this document and classify it for an Obsidian vault.

The vault belongs to a Solution Architect at Elastic. Available folders:
{folder_list}

Valid note types: {type_list}

File context (use these hints where the document content is ambiguous):
{path_context}

Document content (first 4000 chars):
{content[:4000]}

Respond with ONLY valid JSON (no markdown fences):
{{
  "type": "<note type>",
  "title": "<Title Case note title>",
  "target_folder": "<folder name from list above>",
  "frontmatter": {{
    "type": "<note type>",
    ...additional relevant properties with wikilinks for relationships...
  }},
  "confidence": <0.0-1.0>,
  "reasoning": "<one sentence>"
}}

Rules:
- Use wikilinks for customer/project/people references: "[[Name]]"
- Title should be Title Case with spaces
- For dated content, prefix title with YYYY-MM-DD - Title
- If unsure, set target_folder to "00-Inbox" and confidence below 0.5
- Only include frontmatter properties that have actual values
- If a customer or date hint is provided above, use it unless the document clearly contradicts it"""

    message = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=1024,
        messages=[{"role": "user", "content": prompt}],
    )

    return json.loads(message.content[0].text)


# ---------------------------------------------------------------------------
# File writing
# ---------------------------------------------------------------------------

def write_note(
    vault_path: Path,
    target_folder: str,
    title: str,
    frontmatter: dict[str, Any],
    body: str,
) -> Path:
    """Write a note into the vault with frontmatter.

    Appends a numeric suffix if a note with the same title already exists.
    """
    folder_path = vault_path / target_folder
    folder_path.mkdir(parents=True, exist_ok=True)

    safe_title = title.replace("/", "-").replace("\\", "-")
    note_path = folder_path / f"{safe_title}.md"

    # Avoid overwrites by appending a counter
    counter = 1
    while note_path.exists():
        counter += 1
        note_path = folder_path / f"{safe_title} {counter}.md"

    yaml_block = yaml.dump(frontmatter, default_flow_style=False, allow_unicode=True)
    content = f"---\n{yaml_block}---\n# {title}\n\n{body}\n"

    note_path.write_text(content, encoding="utf-8")
    return note_path


# ---------------------------------------------------------------------------
# Main pipeline
# ---------------------------------------------------------------------------

def run_import(
    config_path: Path,
    source_dir: Path,
    dry_run: bool = False,
    use_ai: bool = True,
) -> None:
    """Run the local file import pipeline.

    When *use_ai* is False, classification relies entirely on directory names,
    subdirectory names, and filename pattern matching -- no API calls are made.
    """
    config = load_config(config_path)
    history = load_history(config.tracker_path)

    if not source_dir.exists():
        raise FileNotFoundError(f"Source directory not found: {source_dir}")

    source_files = discover_files(source_dir, config.note_types)
    presorted = sum(1 for f in source_files if f.hint_type)
    unsorted = len(source_files) - presorted

    logger.info("Found %d files (%d with type hint, %d without)", len(source_files), presorted, unsorted)
    if not use_ai:
        logger.info("Running in local-only mode (no AI API calls)")

    if not source_files:
        logger.info("Nothing to import.")
        return

    imported = 0
    skipped = 0
    failed = 0

    for sf in source_files:
        file_key = str(sf.path.resolve())

        if file_key in history:
            logger.info("Skipped (already imported): %s", sf.path.name)
            skipped += 1
            continue

        hints = []
        if sf.hint_type:
            hints.append(f"type={sf.hint_type}")
        if sf.hint_customer:
            hints.append(f"customer={sf.hint_customer}")
        if sf.hint_date:
            hints.append(f"date={sf.hint_date}")
        if sf.hint_archived:
            hints.append("archived")
        hint_label = f" [{', '.join(hints)}]" if hints else " [unclassified]"
        logger.info("Processing: %s%s", sf.path.name, hint_label)

        try:
            markdown_content = convert_file(sf.path)

            if not markdown_content.strip():
                logger.warning("Skipped (empty): %s", sf.path.name)
                skipped += 1
                continue

            if use_ai:
                if sf.hint_type:
                    metadata = extract_metadata(markdown_content, sf, sf.hint_type)
                    title = metadata["title"]
                    frontmatter = metadata["frontmatter"]
                    target_folder = TYPE_TO_FOLDER.get(sf.hint_type, "00-Inbox")
                    note_type = sf.hint_type
                    confidence = 1.0
                    reasoning = f"Pre-sorted into {sf.hint_type}/ directory"
                else:
                    classification = classify_document(
                        markdown_content, sf, config.folders, config.note_types,
                    )
                    title = classification["title"]
                    frontmatter = classification["frontmatter"]
                    target_folder = classification["target_folder"]
                    note_type = classification["type"]
                    confidence = classification.get("confidence", 0.0)
                    reasoning = classification.get("reasoning", "")
            else:
                classification = classify_locally(sf)
                title = classification["title"]
                frontmatter = classification["frontmatter"]
                target_folder = classification["target_folder"]
                note_type = classification["type"]
                confidence = classification.get("confidence", 0.0)
                reasoning = classification.get("reasoning", "")

            # Merge path-derived hints as fallbacks if AI didn't set them
            if sf.hint_date and "date" not in frontmatter:
                frontmatter["date"] = sf.hint_date
            if sf.hint_customer and "customer" not in frontmatter:
                frontmatter["customer"] = f"[[{sf.hint_customer}]]"
            if sf.hint_archived:
                frontmatter["status"] = "archived"

            if dry_run:
                logger.info(
                    "  [DRY RUN] -> %s/%s.md (type=%s, confidence=%.0f%%) -- %s",
                    target_folder, title, note_type, confidence * 100, reasoning,
                )
            else:
                note_path = write_note(
                    config.vault_path, target_folder, title, frontmatter, markdown_content,
                )
                logger.info("  -> %s", note_path.relative_to(config.vault_path))

                history[file_key] = {
                    "source": sf.path.name,
                    "title": title,
                    "folder": target_folder,
                    "type": note_type,
                    "pre_sorted": sf.hint_type is not None,
                }

            imported += 1

        except Exception as exc:
            logger.error("Failed to process %s: %s", sf.path.name, exc)
            failed += 1

    if not dry_run:
        save_history(config.tracker_path, history)

    logger.info("=" * 60)
    logger.info("Import complete: %d processed, %d skipped, %d failed", imported, skipped, failed)
    if dry_run:
        logger.info("DRY RUN -- no files were written")
    logger.info("=" * 60)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> None:
    """CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Import local files (.docx, .pptx, .md, .txt) into your Obsidian vault",
        epilog=(
            "Supports customer-first layouts (CustomerName/files), type-first "
            "layouts (meeting/files), or flat directories. Layout is auto-detected."
        ),
    )
    parser.add_argument(
        "source_dir",
        help="Directory containing files to import (supports subdirectories)",
    )
    parser.add_argument(
        "--config",
        default=str(PROJECT_ROOT / "config.yaml"),
        help="Path to config.yaml (default: ./config.yaml)",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Preview classification without writing files",
    )
    parser.add_argument(
        "--ai",
        action="store_true",
        help="Use Anthropic API for classification (requires ANTHROPIC_API_KEY)",
    )

    args = parser.parse_args()

    run_import(
        config_path=Path(args.config),
        source_dir=Path(args.source_dir).expanduser().resolve(),
        dry_run=args.dry_run,
        use_ai=args.ai,
    )


if __name__ == "__main__":
    main()
